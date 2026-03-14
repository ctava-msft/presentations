"""Test script to verify correct animation XML structure."""
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
from lxml import etree

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Test Fade Animation"
body = slide.shapes.placeholders[1].text_frame
body.clear()
body.paragraphs[0].text = "This should fade in on click"

# Get shape IDs
for shape in slide.shapes:
    sp = shape._element
    nvSpPr = sp.find(qn("p:nvSpPr"))
    if nvSpPr is not None:
        cNvPr = nvSpPr.find(qn("p:cNvPr"))
        print(f"Shape: {cNvPr.get('name')} id={cNvPr.get('id')}")

# Correct animation XML - proper 3-level nesting for click-triggered fade
timing_xml = """<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                          xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
                <p:par>
                  <p:cTn id="3" fill="hold">
                    <p:stCondLst>
                      <p:cond delay="indefinite"/>
                    </p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="4" fill="hold">
                          <p:stCondLst>
                            <p:cond delay="0"/>
                          </p:stCondLst>
                          <p:childTnLst>
                            <p:par>
                              <p:cTn id="5" presetID="10" presetClass="entr" presetSubtype="0" fill="hold" nodeType="clickEffect">
                                <p:stCondLst>
                                  <p:cond delay="0"/>
                                </p:stCondLst>
                                <p:childTnLst>
                                  <p:set>
                                    <p:cBhvr>
                                      <p:cTn id="6" dur="1" fill="hold">
                                        <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                                      </p:cTn>
                                      <p:tgtEl><p:spTgt spid="3"/></p:tgtEl>
                                      <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                                    </p:cBhvr>
                                    <p:to><p:strVal val="visible"/></p:to>
                                  </p:set>
                                  <p:animEffect transition="in" filter="fade">
                                    <p:cBhvr>
                                      <p:cTn id="7" dur="500" fill="hold"/>
                                      <p:tgtEl><p:spTgt spid="3"/></p:tgtEl>
                                    </p:cBhvr>
                                  </p:animEffect>
                                </p:childTnLst>
                              </p:cTn>
                            </p:par>
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst>
              <p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:prevCondLst>
            <p:nextCondLst>
              <p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:nextCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>"""

timing_el = etree.fromstring(timing_xml.encode())
slide._element.append(timing_el)
prs.save("output/test_anim.pptx")
print("Saved output/test_anim.pptx")
