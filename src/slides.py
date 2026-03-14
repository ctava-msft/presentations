"""Slide builder functions – one per layout type.

Each builder receives a ``Presentation``, slide data, a :class:`Style`, and
optional helpers for images and animations.
"""

from __future__ import annotations

import os
from typing import TYPE_CHECKING

from pptx.util import Inches

if TYPE_CHECKING:
    from pptx import Presentation

    from .style import Style

# ---------------------------------------------------------------------------
# Shared helpers
# ---------------------------------------------------------------------------


def _apply_position(shape, pos: dict | None) -> None:
    """Move/resize *shape* according to a parsed position dict (inches)."""
    if not pos:
        return
    if "left" in pos:
        shape.left = Inches(pos["left"])
    if "top" in pos:
        shape.top = Inches(pos["top"])
    if "width" in pos:
        shape.width = Inches(pos["width"])
    if "height" in pos:
        shape.height = Inches(pos["height"])


def _add_image(slide, img: dict, pos: dict | None = None) -> None:
    """Add an image to *slide* from a parsed ``**Image**`` dict.

    If *pos* is given (from ``**ImagePos**``), it overrides the image dict coords.
    """
    path = img["path"]
    if not os.path.isfile(path):
        print(f"Warning: image not found: {path}, skipping.")
        return
    left = Inches((pos or {}).get("left", img.get("left", 6.5)))
    top = Inches((pos or {}).get("top", img.get("top", 1.5)))
    w_val = (pos or {}).get("width", img.get("width"))
    h_val = (pos or {}).get("height", img.get("height"))
    width = Inches(w_val) if w_val is not None else None
    height = Inches(h_val) if h_val is not None else None
    slide.shapes.add_picture(path, left, top, width, height)


# ---------------------------------------------------------------------------
# Layout builders
# ---------------------------------------------------------------------------


def add_title_slide(
    prs: Presentation,
    slide_data: dict,
    style: Style,
    *,
    apply_animations=None,
) -> None:
    """Layout 0 – Title Slide: large centred title + subtitle."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    positions = slide_data.get("positions", {})
    slide.shapes.title.text = slide_data["title"]
    slide.placeholders[1].text = slide_data.get("subtitle", "")
    slide.shapes.title.text_frame.paragraphs[0].font.size = style.title_font
    slide.placeholders[1].text_frame.paragraphs[0].font.size = style.subtitle_font
    _apply_position(slide.shapes.title, positions.get("title"))
    _apply_position(slide.placeholders[1], positions.get("subtitle"))
    slide.notes_slide.notes_text_frame.text = slide_data.get("notes", "")
    if slide_data.get("image"):
        _add_image(slide, slide_data["image"], positions.get("image"))
    if apply_animations and slide_data.get("animations"):
        apply_animations(slide, slide_data["animations"])


def add_content_slide(
    prs: Presentation,
    slide_data: dict,
    style: Style,
    *,
    apply_animations=None,
) -> None:
    """Layout 1 – Title and Content: title bar + bullet list."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    positions = slide_data.get("positions", {})
    slide.shapes.title.text = slide_data["title"]
    slide.shapes.title.text_frame.paragraphs[0].font.size = style.heading_font
    _apply_position(slide.shapes.title, positions.get("title"))
    body_ph = slide.shapes.placeholders[1]
    _apply_position(body_ph, positions.get("content"))
    body = body_ph.text_frame
    body.clear()
    for i, b in enumerate(slide_data.get("bullets", [])):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = style.body_font
    slide.notes_slide.notes_text_frame.text = slide_data.get("notes", "")
    if slide_data.get("image"):
        _add_image(slide, slide_data["image"], positions.get("image"))
    if apply_animations and slide_data.get("animations"):
        apply_animations(slide, slide_data["animations"])


def add_section_header_slide(
    prs: Presentation,
    slide_data: dict,
    style: Style,
    *,
    apply_animations=None,
) -> None:
    """Layout 2 – Section Header: transition slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    positions = slide_data.get("positions", {})
    slide.shapes.title.text = slide_data["title"]
    slide.shapes.title.text_frame.paragraphs[0].font.size = style.title_font
    _apply_position(slide.shapes.title, positions.get("title"))
    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        slide.placeholders[1].text = subtitle
        slide.placeholders[1].text_frame.paragraphs[0].font.size = style.subtitle_font
        _apply_position(slide.placeholders[1], positions.get("subtitle"))
    slide.notes_slide.notes_text_frame.text = slide_data.get("notes", "")
    if slide_data.get("image"):
        _add_image(slide, slide_data["image"], positions.get("image"))
    if apply_animations and slide_data.get("animations"):
        apply_animations(slide, slide_data["animations"])


def add_two_column_slide(
    prs: Presentation,
    slide_data: dict,
    style: Style,
    *,
    apply_animations=None,
) -> None:
    """Layout 3 – Two Content: side-by-side content placeholders."""
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    positions = slide_data.get("positions", {})
    slide.shapes.title.text = slide_data["title"]
    slide.shapes.title.text_frame.paragraphs[0].font.size = style.heading_font
    _apply_position(slide.shapes.title, positions.get("title"))

    left_placeholder = slide.placeholders[1]
    _apply_position(left_placeholder, positions.get("left"))
    left_ph = left_placeholder.text_frame
    left_ph.clear()
    for i, b in enumerate(slide_data.get("left_bullets", [])):
        p = left_ph.paragraphs[0] if i == 0 else left_ph.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = style.col_body_font

    right_placeholder = slide.placeholders[2]
    _apply_position(right_placeholder, positions.get("right"))
    right_ph = right_placeholder.text_frame
    right_ph.clear()
    for i, b in enumerate(slide_data.get("right_bullets", [])):
        p = right_ph.paragraphs[0] if i == 0 else right_ph.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = style.col_body_font

    slide.notes_slide.notes_text_frame.text = slide_data.get("notes", "")
    if slide_data.get("image"):
        _add_image(slide, slide_data["image"], positions.get("image"))
    if apply_animations and slide_data.get("animations"):
        apply_animations(slide, slide_data["animations"])


# ---------------------------------------------------------------------------
# Registry mapping slide type names → builder functions
# ---------------------------------------------------------------------------

SLIDE_BUILDERS: dict[str, callable] = {
    "title": add_title_slide,
    "content": add_content_slide,
    "section-header": add_section_header_slide,
    "two-column": add_two_column_slide,
}
