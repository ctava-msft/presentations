"""
Presentation generator: reads a .spec.md file and produces a PowerPoint deck.

Usage:
    python presentation.py <spec-file>

Example:
    python presentation.py .speckit/specifications/ai101.spec.md
"""

import argparse
import os
import re
import sys
from copy import deepcopy
from lxml import etree

import yaml
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn, nsmap
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Theme constants
# ---------------------------------------------------------------------------
TITLE_FONT = Pt(36)
SUBTITLE_FONT = Pt(20)
BODY_FONT = Pt(20)
HEADING_FONT = Pt(32)
COL_HEADING_FONT = Pt(22)
COL_BODY_FONT = Pt(18)

# ---------------------------------------------------------------------------
# Spec parser
# ---------------------------------------------------------------------------

def parse_spec(path: str) -> dict:
    """Parse a presentation spec markdown file into metadata + slide list."""
    with open(path, encoding="utf-8") as f:
        text = f.read()

    # Split YAML front matter
    fm_match = re.match(r"^---\s*\n(.*?)\n---\s*\n", text, re.DOTALL)
    if not fm_match:
        sys.exit("Error: spec file must start with YAML front matter (--- … ---)")
    metadata = yaml.safe_load(fm_match.group(1))
    body = text[fm_match.end():]

    # Split slides on horizontal rules (--- on its own line)
    raw_slides = re.split(r"\n---\s*\n", body)
    slides = []
    for raw in raw_slides:
        raw = raw.strip()
        if not raw:
            continue
        slide = _parse_slide(raw)
        if slide:
            slides.append(slide)

    return {"metadata": metadata, "slides": slides}


def _parse_slide(raw: str) -> dict | None:
    """Parse a single slide block into a structured dict."""
    # Header: ## [type] Title
    header_match = re.match(r"^##\s+\[(\w[\w-]*)\]\s+(.+)", raw)
    if not header_match:
        return None
    slide_type = header_match.group(1).strip()
    title = header_match.group(2).strip()
    rest = raw[header_match.end():].strip()

    # Extract notes (everything after **Notes**: to end)
    notes = ""
    notes_match = re.split(r"\*\*Notes\*\*\s*:\s*", rest, maxsplit=1)
    if len(notes_match) == 2:
        rest = notes_match[0].strip()
        notes = notes_match[1].strip()

    slide: dict = {"type": slide_type, "title": title, "notes": notes}

    # --- Image (optional, any slide type) ---
    slide["image"] = _parse_image_field(rest)

    # --- Animations (optional, any slide type) ---
    slide["animations"] = _parse_animations(rest)

    if slide_type == "title":
        sub_match = re.search(r"\*\*Subtitle\*\*\s*:\s*(.+)", rest)
        slide["subtitle"] = sub_match.group(1).strip() if sub_match else ""

    elif slide_type == "section-header":
        sub_match = re.search(r"\*\*Subtitle\*\*\s*:\s*(.+)", rest)
        slide["subtitle"] = sub_match.group(1).strip() if sub_match else ""

    elif slide_type == "content":
        slide["bullets"] = _extract_bullets(rest)

    elif slide_type == "two-column":
        slide.update(_parse_two_column(rest))

    return slide


def _extract_bullets(text: str) -> list[str]:
    """Extract markdown list items from text."""
    return [m.group(1).strip() for m in re.finditer(r"^[-*]\s+(.+)", text, re.MULTILINE)]


def _parse_two_column(text: str) -> dict:
    """Parse two-column slide fields."""
    result: dict = {}

    # Extract left bullets (between **Left**: and **Right**: or **Notes**: or end)
    left_match = re.search(
        r"\*\*Left\*\*\s*:\s*\n(.*?)(?=\*\*Right\*\*|\*\*Notes\*\*|\*\*Image\*\*|\*\*Animation\*\*|$)",
        text, re.DOTALL,
    )
    result["left_bullets"] = _extract_bullets(left_match.group(1)) if left_match else []

    # Extract right bullets (between **Right**: and **Notes**: or end)
    right_match = re.search(
        r"\*\*Right\*\*\s*:\s*\n(.*?)(?=\*\*Notes\*\*|\*\*Image\*\*|\*\*Animation\*\*|$)",
        text, re.DOTALL,
    )
    result["right_bullets"] = _extract_bullets(right_match.group(1)) if right_match else []

    return result


def _parse_image_field(text: str) -> dict | None:
    """Parse optional **Image**: path, position, size."""
    match = re.search(r"\*\*Image\*\*\s*:\s*(.+)", text)
    if not match:
        return None
    raw = match.group(1).strip()
    # Format: path [, left, top, width, height]  (inches, all optional after path)
    parts = [p.strip() for p in raw.split(",")]
    img: dict = {"path": parts[0]}
    if len(parts) >= 3:
        img["left"] = float(parts[1])
        img["top"] = float(parts[2])
    if len(parts) >= 5:
        img["width"] = float(parts[3])
        img["height"] = float(parts[4])
    return img


def _parse_animations(text: str) -> list[dict]:
    """Parse **Animation**: lines.  Each line: target > effect [, options]."""
    animations = []
    for m in re.finditer(r"\*\*Animation\*\*\s*:\s*(.+)", text):
        raw = m.group(1).strip()
        # Format:  target > effect
        #   target = title | content | image | left | right | bullets | all
        #   effect = appear | fade | fly-in | wipe | zoom | float-in
        parts = [p.strip() for p in raw.split(">", 1)]
        if len(parts) == 2:
            target, effect_str = parts
        else:
            target, effect_str = "all", parts[0]
        animations.append({"target": target.lower(), "effect": effect_str.lower().strip()})
    return animations


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def add_title_slide(prs: Presentation, title: str, subtitle: str, notes: str,
                    image: dict | None = None, animations: list[dict] | None = None):
    """Layout 0 – Title Slide: large centered title + subtitle."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT
    slide.placeholders[1].text_frame.paragraphs[0].font.size = SUBTITLE_FONT
    slide.notes_slide.notes_text_frame.text = notes
    if image:
        _add_image(slide, image)
    if animations:
        _apply_animations(slide, animations)


def add_content_slide(prs: Presentation, title: str, bullets: list[str], notes: str,
                      image: dict | None = None, animations: list[dict] | None = None):
    """Layout 1 – Title and Content: title bar + single content placeholder."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = HEADING_FONT
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = BODY_FONT
    slide.notes_slide.notes_text_frame.text = notes
    if image:
        _add_image(slide, image)
    if animations:
        _apply_animations(slide, animations)


def add_section_header_slide(prs: Presentation, title: str, subtitle: str, notes: str,
                             image: dict | None = None, animations: list[dict] | None = None):
    """Layout 2 – Section Header: large title + subtitle for topic transitions."""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT
    if subtitle:
        slide.placeholders[1].text = subtitle
        slide.placeholders[1].text_frame.paragraphs[0].font.size = SUBTITLE_FONT
    slide.notes_slide.notes_text_frame.text = notes
    if image:
        _add_image(slide, image)
    if animations:
        _apply_animations(slide, animations)


def add_two_column_slide(
    prs: Presentation,
    title: str,
    left_bullets: list[str],
    right_bullets: list[str],
    notes: str,
    image: dict | None = None,
    animations: list[dict] | None = None,
):
    """Layout 3 – Two Content: title + two side-by-side content placeholders."""
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = HEADING_FONT

    # Left content – placeholder index 1
    left_ph = slide.placeholders[1].text_frame
    left_ph.clear()
    for i, b in enumerate(left_bullets):
        p = left_ph.paragraphs[0] if i == 0 else left_ph.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = COL_BODY_FONT

    # Right content – placeholder index 2
    right_ph = slide.placeholders[2].text_frame
    right_ph.clear()
    for i, b in enumerate(right_bullets):
        p = right_ph.paragraphs[0] if i == 0 else right_ph.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = COL_BODY_FONT

    slide.notes_slide.notes_text_frame.text = notes
    if image:
        _add_image(slide, image)
    if animations:
        _apply_animations(slide, animations)


# ---------------------------------------------------------------------------
# Image helper
# ---------------------------------------------------------------------------

def _add_image(slide, img: dict):
    """Add an image to the slide from the spec's **Image** field."""
    path = img["path"]
    if not os.path.isfile(path):
        print(f"Warning: image not found: {path}, skipping.")
        return
    left = Inches(img.get("left", 6.5))
    top = Inches(img.get("top", 1.5))
    width = Inches(img.get("width", 3.0)) if "width" in img else None
    height = Inches(img.get("height", 3.0)) if "height" in img else None
    slide.shapes.add_picture(path, left, top, width, height)


# ---------------------------------------------------------------------------
# Animation engine  (direct Open XML injection)
# ---------------------------------------------------------------------------

# Map of plain-English effect names → PowerPoint preset + visual effect config
# Each entry: preset ID, preset class, subtype, and an effect-builder key
_EFFECT_MAP = {
    "appear":       {"preset": "1",  "cls": "entr", "subtype": "0",  "visual": "none"},
    "fade":         {"preset": "10", "cls": "entr", "subtype": "0",  "visual": "fade"},
    "fly-in":       {"preset": "2",  "cls": "entr", "subtype": "4",  "visual": "fly-bottom"},
    "fly-in-left":  {"preset": "2",  "cls": "entr", "subtype": "8",  "visual": "fly-left"},
    "fly-in-right": {"preset": "2",  "cls": "entr", "subtype": "2",  "visual": "fly-right"},
    "fly-in-top":   {"preset": "2",  "cls": "entr", "subtype": "1",  "visual": "fly-top"},
    "wipe":         {"preset": "22", "cls": "entr", "subtype": "4",  "visual": "wipe"},
    "zoom":         {"preset": "23", "cls": "entr", "subtype": "0",  "visual": "zoom"},
    "float-in":     {"preset": "42", "cls": "entr", "subtype": "4",  "visual": "float"},
    "split":        {"preset": "16", "cls": "entr", "subtype": "0",  "visual": "split"},
    "blinds":       {"preset": "3",  "cls": "entr", "subtype": "0",  "visual": "blinds"},
}

# Map of target keywords → shape name patterns
_TARGET_SHAPE_MAP = {
    "title":    ["Title"],
    "content":  ["Content Placeholder", "Text Placeholder"],
    "left":     ["Content Placeholder 2"],
    "right":    ["Content Placeholder 3"],
    "image":    ["Picture"],
    "subtitle": ["Subtitle"],
}

_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _resolve_target_shapes(slide, target: str) -> list:
    """Resolve an English target keyword to actual shape objects on the slide."""
    if target == "all":
        return list(slide.shapes)
    if target == "bullets":
        target = "content"
    patterns = _TARGET_SHAPE_MAP.get(target, [])
    matched = []
    for shape in slide.shapes:
        for pat in patterns:
            if pat.lower() in shape.name.lower():
                matched.append(shape)
                break
    if not matched:
        for shape in slide.shapes:
            if target.lower() in shape.name.lower():
                matched.append(shape)
    return matched


def _get_shape_id(shape) -> str:
    """Get the numeric id attribute from a shape's XML."""
    sp_elem = shape._element
    for nv_tag in ("p:nvSpPr", "p:nvPicPr", "p:nvGrpSpPr", "p:nvCxnSpPr"):
        nv = sp_elem.find(qn(nv_tag))
        if nv is not None:
            cNvPr = nv.find(qn("p:cNvPr"))
            if cNvPr is not None:
                return cNvPr.get("id")
    return "2"


def _build_visual_effect_xml(shape_id: str, visual: str) -> str:
    """Return extra animation XML nodes for the specific visual effect."""
    if visual == "none":
        # Appear: no extra nodes needed beyond the <p:set>
        return ""
    elif visual == "fade":
        return f"""<p:animEffect transition="in" filter="fade">
            <p:cBhvr>
              <p:cTn id="0" dur="500" fill="hold"/>
              <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
            </p:cBhvr>
          </p:animEffect>"""
    elif visual.startswith("fly-"):
        # Fly-in uses <p:anim> to animate ppt_x or ppt_y
        direction = visual.split("-", 1)[1]
        if direction == "bottom":
            attr, fr, to = "ppt_y", "#ppt_h+#ppt_y", "#ppt_y"
        elif direction == "top":
            attr, fr, to = "ppt_y", "-#ppt_h", "#ppt_y"
        elif direction == "left":
            attr, fr, to = "ppt_x", "-#ppt_w", "#ppt_x"
        elif direction == "right":
            attr, fr, to = "ppt_x", "#ppt_w+#ppt_x", "#ppt_x"
        else:
            attr, fr, to = "ppt_y", "#ppt_h+#ppt_y", "#ppt_y"
        return f"""<p:anim calcmode="lin" valueType="num">
            <p:cBhvr additive="base">
              <p:cTn id="0" dur="500" fill="hold"/>
              <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
              <p:attrNameLst><p:attrName>{attr}</p:attrName></p:attrNameLst>
            </p:cBhvr>
            <p:tavLst>
              <p:tav tm="0"><p:val><p:strVal val="{fr}"/></p:val></p:tav>
              <p:tav tm="100000"><p:val><p:strVal val="{to}"/></p:val></p:tav>
            </p:tavLst>
          </p:anim>"""
    elif visual == "wipe":
        return f"""<p:animEffect transition="in" filter="wipe(down)">
            <p:cBhvr>
              <p:cTn id="0" dur="500" fill="hold"/>
              <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
            </p:cBhvr>
          </p:animEffect>"""
    elif visual == "zoom":
        return f"""<p:animScale>
            <p:cBhvr>
              <p:cTn id="0" dur="500" fill="hold"/>
              <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
            </p:cBhvr>
            <p:by x="0" y="0"/>
            <p:from x="0" y="0"/>
            <p:to x="100000" y="100000"/>
          </p:animScale>"""
    elif visual == "float":
        return f"""<p:anim calcmode="lin" valueType="num">
            <p:cBhvr additive="base">
              <p:cTn id="0" dur="500" fill="hold"/>
              <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
              <p:attrNameLst><p:attrName>ppt_y</p:attrName></p:attrNameLst>
            </p:cBhvr>
            <p:tavLst>
              <p:tav tm="0"><p:val><p:strVal val="#ppt_y+0.1"/></p:val></p:tav>
              <p:tav tm="100000"><p:val><p:strVal val="#ppt_y"/></p:val></p:tav>
            </p:tavLst>
          </p:anim>
          <p:animEffect transition="in" filter="fade">
            <p:cBhvr>
              <p:cTn id="0" dur="500" fill="hold"/>
              <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
            </p:cBhvr>
          </p:animEffect>"""
    elif visual == "split":
        return f"""<p:animEffect transition="in" filter="barn(inVertical)">
            <p:cBhvr>
              <p:cTn id="0" dur="500" fill="hold"/>
              <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
            </p:cBhvr>
          </p:animEffect>"""
    elif visual == "blinds":
        return f"""<p:animEffect transition="in" filter="blinds(horizontal)">
            <p:cBhvr>
              <p:cTn id="0" dur="500" fill="hold"/>
              <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
            </p:cBhvr>
          </p:animEffect>"""
    return ""


def _build_click_par(shape_id: str, effect: dict) -> str:
    """Build one complete click-triggered animation group (3-level <p:par> nesting)."""
    preset = effect["preset"]
    cls = effect["cls"]
    subtype = effect["subtype"]
    visual = effect["visual"]
    visual_xml = _build_visual_effect_xml(shape_id, visual)

    return f"""<p:par xmlns:p="{_P_NS}" xmlns:a="{_A_NS}">
  <p:cTn id="0" fill="hold">
    <p:stCondLst>
      <p:cond delay="indefinite"/>
    </p:stCondLst>
    <p:childTnLst>
      <p:par>
        <p:cTn id="0" fill="hold">
          <p:stCondLst>
            <p:cond delay="0"/>
          </p:stCondLst>
          <p:childTnLst>
            <p:par>
              <p:cTn id="0" presetID="{preset}" presetClass="{cls}" presetSubtype="{subtype}"
                     fill="hold" nodeType="clickEffect">
                <p:stCondLst>
                  <p:cond delay="0"/>
                </p:stCondLst>
                <p:childTnLst>
                  <p:set>
                    <p:cBhvr>
                      <p:cTn id="0" dur="1" fill="hold">
                        <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                      </p:cTn>
                      <p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>
                      <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                    </p:cBhvr>
                    <p:to><p:strVal val="visible"/></p:to>
                  </p:set>
                  {visual_xml}
                </p:childTnLst>
              </p:cTn>
            </p:par>
          </p:childTnLst>
        </p:cTn>
      </p:par>
    </p:childTnLst>
  </p:cTn>
</p:par>"""


def _apply_animations(slide, animations: list[dict]):
    """Inject animation XML into the slide for each animation spec."""
    if not animations:
        return

    # Collect all (shape_id, effect_dict) pairs
    anim_pairs = []
    for anim in animations:
        effect_name = anim["effect"]
        effect = _EFFECT_MAP.get(effect_name)
        if effect is None:
            print(f"Warning: unknown animation '{effect_name}', skipping. "
                  f"Available: {', '.join(_EFFECT_MAP.keys())}")
            continue
        shapes = _resolve_target_shapes(slide, anim["target"])
        if not shapes:
            print(f"Warning: no shapes matched target '{anim['target']}' on slide, skipping animation.")
            continue
        for s in shapes:
            anim_pairs.append((_get_shape_id(s), effect))

    if not anim_pairs:
        return

    # Build click-par blocks for each animation
    click_pars = ""
    for shape_id, eff in anim_pairs:
        click_pars += _build_click_par(shape_id, eff)

    timing_xml = f"""<p:timing xmlns:p="{_P_NS}" xmlns:a="{_A_NS}">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
                {click_pars}
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst>
              <p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:prevCondLst>
            <p:nextCondLst>
              <p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:nextCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>"""

    timing_el = etree.fromstring(timing_xml.encode())

    # Assign unique sequential IDs to all cTn elements
    ctn_id = 1
    for ctn in timing_el.iter(qn("p:cTn")):
        ctn.set("id", str(ctn_id))
        ctn_id += 1

    # Remove any existing <p:timing> and append ours
    existing = slide._element.find(qn("p:timing"))
    if existing is not None:
        slide._element.remove(existing)
    slide._element.append(timing_el)


# ---------------------------------------------------------------------------
# Renderer
# ---------------------------------------------------------------------------

SLIDE_BUILDERS = {
    "title": lambda prs, s: add_title_slide(
        prs, s["title"], s.get("subtitle", ""), s.get("notes", ""),
        s.get("image"), s.get("animations")),
    "content": lambda prs, s: add_content_slide(
        prs, s["title"], s.get("bullets", []), s.get("notes", ""),
        s.get("image"), s.get("animations")),
    "section-header": lambda prs, s: add_section_header_slide(
        prs, s["title"], s.get("subtitle", ""), s.get("notes", ""),
        s.get("image"), s.get("animations")),
    "two-column": lambda prs, s: add_two_column_slide(
        prs, s["title"], s.get("left_bullets", []), s.get("right_bullets", []),
        s.get("notes", ""), s.get("image"), s.get("animations")),
}


def render(spec: dict, output_dir: str = "output"):
    """Render a parsed spec into a PowerPoint file."""
    metadata = spec["metadata"]
    slides = spec["slides"]
    out_name = metadata.get("output", "presentation.pptx")

    prs = Presentation()

    for slide_data in slides:
        stype = slide_data["type"]
        builder = SLIDE_BUILDERS.get(stype)
        if builder is None:
            print(f"Warning: unknown slide type '{stype}', skipping.")
            continue
        builder(prs, slide_data)

    os.makedirs(output_dir, exist_ok=True)
    out_path = _next_version_path(output_dir, out_name)
    prs.save(out_path)
    print(f"Saved {len(slides)} slides -> {out_path}")
    return out_path


def _next_version_path(output_dir: str, filename: str) -> str:
    """Return a versioned path: filename.pptx, filename_1.pptx, filename_2.pptx, …"""
    base, ext = os.path.splitext(filename)
    candidate = os.path.join(output_dir, filename)
    n = 1
    while os.path.exists(candidate):
        candidate = os.path.join(output_dir, f"{base}_{n}{ext}")
        n += 1
    return candidate


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="Generate a PowerPoint presentation from a spec file.")
    parser.add_argument("spec", help="Path to the .spec.md file")
    parser.add_argument("-o", "--output-dir", default="output", help="Output directory (default: output)")
    args = parser.parse_args()

    if not os.path.isfile(args.spec):
        sys.exit(f"Error: spec file not found: {args.spec}")

    spec = parse_spec(args.spec)
    render(spec, args.output_dir)


if __name__ == "__main__":
    main()
